from flask import Flask, render_template, request
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os

app = Flask(__name__)

EXCEL_FILE = "mythos.xlsx"
PPT_FILE = "mythos_assessment_report.pptx"

QUESTIONS = [
    {"domain": "Identity inventory", "question": "Do you have a complete inventory of all human, service, machine, workload, and AI-agent identities?"},
    {"domain": "AI agent identity", "question": "Can you distinguish between a human user, automation script, service account, and AI agent in logs?"},
    {"domain": "Authentication", "question": "Which IdP is used: Entra ID, Okta, Ping, Oracle IAM, or others?"},
    {"domain": "MFA", "question": "Is MFA enforced for all privileged users, admins, contractors, and service owners?"},
    {"domain": "Conditional Access", "question": "Do you enforce device, location, risk, session, and workload-based access policies?"},
    {"domain": "Privileged access", "question": "Do admins use permanent access or just-in-time access through PIM/PAM?"},
    {"domain": "Service accounts", "question": "How many service accounts exist, and are any using static passwords or long-lived secrets?"},
    {"domain": "API access", "question": "Are APIs protected by OAuth, mTLS, API keys, APIM, gateway policies, or network restrictions?"},
    {"domain": "Token controls", "question": "Are JWTs validated for issuer, audience, expiry, scopes, roles, and signature?"},
    {"domain": "Secrets management", "question": "Are API keys, passwords, certificates, and model credentials stored in Key Vault or a secrets manager?"},
    {"domain": "Agent permissions", "question": "What systems can AI agents access: email, HR, CRM, code repos, cloud, tickets, databases, or identity systems?"},
    {"domain": "Least privilege", "question": "Are AI agents restricted to minimum required permissions?"},
    {"domain": "Human approval", "question": "Which AI-agent actions require human approval before execution?"},
    {"domain": "Tool access", "question": "Are AI agents allowed to call external tools, APIs, browsers, shells, or code execution environments?"},
    {"domain": "Policy enforcement", "question": "Do you use OPA, PBAC, ABAC, IAM policies, or workflow approval gates before agent actions?"},
    {"domain": "SaaS identity", "question": "Which SaaS platforms are connected to SSO, and which still use local accounts?"},
    {"domain": "Joiner-mover-leaver", "question": "How quickly are user and contractor accounts disabled after role change or exit?"},
    {"domain": "Identity governance", "question": "Do you run access reviews, SoD checks, entitlement reviews, and orphan-account cleanup?"},
    {"domain": "Logging", "question": "Do you log identity events from IdP, PAM, cloud, SaaS, APIM, and AI-agent platforms?"},
    {"domain": "Detection", "question": "Can you detect abnormal login, impossible travel, token replay, API abuse, or high-speed scanning?"},
    {"domain": "AI-agent behavior", "question": "Can you detect non-human behavior such as rapid enumeration, repeated API probing, or automated privilege testing?"},
    {"domain": "Model inventory", "question": "Do you know which AI models are deployed, where they run, and who can access them?"},
    {"domain": "Shadow AI", "question": "Can users deploy or access unauthorized AI models or agents?"},
    {"domain": "Network segmentation", "question": "Are identity systems, APIs, databases, and AI platforms microsegmented?"},
    {"domain": "Data access", "question": "Are retrieval systems protected with tenant, role, sensitivity, and data-label filters?"},
    {"domain": "Incident response", "question": "Do you have a playbook for compromised service account, AI agent misuse, or rogue model access?"},
    {"domain": "Mythos readiness", "question": "If an AI model finds identity vulnerabilities rapidly, can you prioritize, contain, and remediate before exploitation?"},
    {"domain": "Compensating controls", "question": "Where patching is delayed, do you have containment controls such as MFA, CA, PAM, segmentation, APIM throttling, and monitoring?"},
    {"domain": "Third-party access", "question": "Do vendors, contractors, and external AI tools have restricted and monitored access?"},
    {"domain": "Evidence", "question": "Can you produce audit evidence showing who accessed what, through which identity, and whether it was human or AI-driven?"}
]


def create_excel():
    if not os.path.exists(EXCEL_FILE):
        wb = Workbook()
        ws = wb.active
        ws.title = "Assessment"
        ws.append(["Timestamp", "Domain", "Question", "Answer", "Risk"])
        wb.save(EXCEL_FILE)


def create_powerpoint(rows, risk_percentage, yes_count, no_count, tbd_count):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.6))
    title.text_frame.text = "Mythos Identity Assessment Report"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True

    risk_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.6))
    risk_box.text_frame.text = (
        f"Overall Mythos Risk: {risk_percentage}% | "
        f"Yes: {yes_count}, No: {no_count}, TBD: {tbd_count} "
        f"(TBD counted as No for %)"
    )
    risk_box.text_frame.paragraphs[0].font.size = Pt(16)

    rows_per_slide = 3

    for i in range(0, len(rows), rows_per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.4))
        header.text_frame.text = "Assessment Results"
        header.text_frame.paragraphs[0].font.size = Pt(20)
        header.text_frame.paragraphs[0].font.bold = True

        risk_text = slide.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.3))
        risk_text.text_frame.text = (
            f"Mythos Risk: {risk_percentage}% | "
            f"Yes: {yes_count}, No: {no_count}, TBD: {tbd_count}"
        )
        risk_text.text_frame.paragraphs[0].font.size = Pt(12)

        chunk = rows[i:i + rows_per_slide]

        table = slide.shapes.add_table(
            rows=len(chunk) + 1,
            cols=4,
            left=Inches(0.35),
            top=Inches(1.1),
            width=Inches(12.6),
            height=Inches(5.7)
        ).table

        table.columns[0].width = Inches(1.8)
        table.columns[1].width = Inches(4.2)
        table.columns[2].width = Inches(5.0)
        table.columns[3].width = Inches(1.1)

        headers = ["Domain", "Question", "Answer", "Risk"]

        # Header row formatting
        for col, h in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = h

            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(161, 0, 255)  # #a100ff

            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Body rows
        for r_idx, item in enumerate(chunk, start=1):
            values = [
                item["domain"],
                item["question"],
                item["answer"],
                item["risk"]
            ]

            for c_idx, val in enumerate(values):
                cell = table.cell(r_idx, c_idx)
                cell.text = val

                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(PPT_FILE)


@app.route("/", methods=["GET", "POST"])
def index():
    message = ""

    if request.method == "POST":
        create_excel()

        wb = load_workbook(EXCEL_FILE)
        ws = wb["Assessment"]

        rows = []
        yes_count = 0
        no_count = 0
        tbd_count = 0

        for item in QUESTIONS:
            domain = item["domain"]
            question = item["question"]
            answer = request.form.get(domain, "")
            risk = request.form.get(domain + "_risk", "")

            if risk == "Yes":
                yes_count += 1
            elif risk == "No":
                no_count += 1
            elif risk == "TBD":
                tbd_count += 1
                no_count += 1  # TBD counted as No for risk %

            rows.append({
                "domain": domain,
                "question": question,
                "answer": answer,
                "risk": risk
            })

            ws.append([
                datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                domain,
                question,
                answer,
                risk
            ])

        total = yes_count + no_count
        risk_percentage = round((yes_count / total) * 100, 2) if total > 0 else 0

        wb.save(EXCEL_FILE)
        create_powerpoint(rows, risk_percentage, yes_count, no_count, tbd_count)

        message = (
            f"Saved. PPT generated. Mythos Risk: {risk_percentage}% "
            f"| Yes: {yes_count}, No: {no_count - tbd_count}, TBD: {tbd_count}"
        )

    return render_template("index.html", questions=QUESTIONS, message=message)


if __name__ == "__main__":
    create_excel()
    app.run(debug=True)