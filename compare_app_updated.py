from flask import Flask, render_template, request, url_for
import pandas as pd
import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


app = Flask(__name__)

GOLD_FILE = r"C:\mythos\Gold-Standard\mythos_gold_standard.xlsx"
REPORT_FOLDER = "static"
PPT_FILE_NAME = "mythos_comparison_report.pptx"
PPT_FILE_PATH = os.path.join(REPORT_FOLDER, PPT_FILE_NAME)


def clean_text(text):
    if pd.isna(text):
        return ""

    text = str(text).lower()
    text = re.sub(r"[^a-z0-9\s]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def normalize_domain(text):
    return clean_text(text)


def compare_answers(client_answer, gold_answer):
    client = clean_text(client_answer)
    gold = clean_text(gold_answer)

    if not client:
        return "High", 0

    if not gold:
        return "High", 0

    stop_words = {
        "with", "that", "this", "from", "they", "have", "into",
        "such", "using", "all", "and", "the", "for", "are", "can",
        "you", "your", "what", "which", "how", "does", "been",
        "will", "must", "should", "also", "only", "than", "then",
        "their", "there", "where", "when", "who", "why"
    }

    gold_words = set(gold.split())
    client_words = set(client.split())

    important_words = {
        word for word in gold_words
        if len(word) > 3 and word not in stop_words
    }

    if not important_words:
        return "High", 0

    matched_words = important_words.intersection(client_words)
    score = len(matched_words) / len(important_words)
    score_percent = round(score * 100, 2)

    if score >= 0.65:
        return "Low", score_percent
    elif score >= 0.30:
        return "Medium", score_percent
    else:
        return "High", score_percent


def load_client_file(uploaded_file):
    df = pd.read_excel(uploaded_file, header=None)

    # Expected client Excel format:
    # 0 = Timestamp
    # 1 = Domain
    # 2 = Question
    # 3 = Answer
    # 4 = Risk selected in app.py: Yes / No / TBD
    df = df.iloc[:, :5]
    df.columns = ["Timestamp", "Domain", "Question", "Answer", "SelectedRisk"]

    # Remove accidental header row if present
    df = df[df["Domain"].astype(str).str.lower() != "domain"]

    return df


def load_gold_file():
    gold_df = pd.read_excel(GOLD_FILE)
    gold_df.columns = [str(col).strip() for col in gold_df.columns]

    required_columns = {"Domain", "Gold Standard Answer"}
    missing_columns = required_columns - set(gold_df.columns)

    if missing_columns:
        raise ValueError(
            "Gold standard file is missing required column(s): "
            + ", ".join(missing_columns)
        )

    gold_map = {}

    for _, row in gold_df.iterrows():
        domain = normalize_domain(row["Domain"])
        answer = row["Gold Standard Answer"]
        gold_map[domain] = answer

    return gold_map


def classify_overall_risk(risk_percentage):
    if risk_percentage <= 20:
        return "Low Risk"
    elif risk_percentage <= 50:
        return "Medium Risk"
    else:
        return "High Risk"


def calculate_summary(results):
    total = len(results)

    low_count = sum(1 for r in results if r["risk_level"] == "Low")
    medium_count = sum(1 for r in results if r["risk_level"] == "Medium")
    high_count = sum(1 for r in results if r["risk_level"] == "High")

    yes_count = sum(1 for r in results if str(r["selected_risk"]).strip().lower() == "yes")
    no_count = sum(1 for r in results if str(r["selected_risk"]).strip().lower() == "no")
    tbd_count = sum(1 for r in results if str(r["selected_risk"]).strip().lower() == "tbd")

    yes_percent = round((yes_count / total) * 100, 2) if total else 0
    no_percent = round((no_count / total) * 100, 2) if total else 0
    tbd_percent = round((tbd_count / total) * 100, 2) if total else 0

    low_percent = round((low_count / total) * 100, 2) if total else 0
    medium_percent = round((medium_count / total) * 100, 2) if total else 0
    high_percent = round((high_count / total) * 100, 2) if total else 0

    # Requested rule:
    # If client selected TBD, add it under Yes Risk.
    yes_risk_count = no_count + tbd_count
    yes_risk_percent = round((yes_risk_count / total) * 100, 2) if total else 0

    # Overall risk based on actual gold-standard comparison.
    # Weighted model: High = 1.0, Medium = 0.5, Low = 0.
    weighted_risk_score = high_count + (medium_count * 0.5)
    comparison_risk_percent = round((weighted_risk_score / total) * 100, 2) if total else 0
    overall_risk_level = classify_overall_risk(comparison_risk_percent)

    return {
        "total": total,

        "low": low_count,
        "medium": medium_count,
        "high": high_count,

        "low_percent": low_percent,
        "medium_percent": medium_percent,
        "high_percent": high_percent,

        "yes": yes_count,
        "no": no_count,
        "tbd": tbd_count,

        "yes_percent": yes_percent,
        "no_percent": no_percent,
        "tbd_percent": tbd_percent,

        "yes_risk_count": yes_risk_count,
        "yes_risk_percent": yes_risk_percent,

        "comparison_risk_percent": comparison_risk_percent,
        "overall_risk_level": overall_risk_level
    }


def set_cell_text(cell, text, font_size=9, bold=False, font_color=RGBColor(0, 0, 0)):
    cell.text = str(text) if text is not None else ""

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = font_color


def add_summary_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.4), Inches(0.5))
    title.text_frame.text = "Mythos Gold Standard Comparison Report"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(161, 0, 255)

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.4), Inches(0.35))
    subtitle.text_frame.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    subtitle.text_frame.paragraphs[0].font.size = Pt(12)

    risk_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(12.4), Inches(0.6))
    risk_box.text_frame.text = (
        f"Overall Comparison Risk: {summary['overall_risk_level']} "
        f"({summary['comparison_risk_percent']}%)"
    )
    risk_box.text_frame.paragraphs[0].font.size = Pt(20)
    risk_box.text_frame.paragraphs[0].font.bold = True

    table = slide.shapes.add_table(
        rows=8,
        cols=3,
        left=Inches(0.7),
        top=Inches(2.25),
        width=Inches(11.9),
        height=Inches(4.3)
    ).table

    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)

    headers = ["Metric", "Count", "Percentage"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        set_cell_text(cell, header, font_size=12, bold=True, font_color=RGBColor(255, 255, 255))
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(161, 0, 255)

    data = [
        ("Total Questions Compared", summary["total"], "100%"),
        ("Gold Match - Low Risk", summary["low"], f"{summary['low_percent']}%"),
        ("Gold Partial - Medium Risk", summary["medium"], f"{summary['medium_percent']}%"),
        ("Gold Mismatch - High Risk", summary["high"], f"{summary['high_percent']}%"),
        ("Client Selected Yes", summary["yes"], f"{summary['yes_percent']}%"),
        ("Client Selected No", summary["no"], f"{summary['no_percent']}%"),
        ("Client Selected TBD", summary["tbd"], f"{summary['tbd_percent']}%"),
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            set_cell_text(table.cell(row_idx, col_idx), value, font_size=10)


def add_logic_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.4), Inches(0.5))
    title.text_frame.text = "Risk Logic Used"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.5))
    body.text_frame.text = (
        "1. Each client answer is compared against the Gold Standard Answer.\n"
        "2. Match score >= 65% = Low Risk.\n"
        "3. Match score >= 30% and < 65% = Medium Risk.\n"
        "4. Match score < 30% = High Risk.\n\n"
        "Overall Comparison Risk Formula:\n"
        "Risk % = ((High Count × 1.0) + (Medium Count × 0.5)) / Total × 100\n\n"
        "Low Risk: 0% - 20%\n"
        "Medium Risk: 21% - 50%\n"
        "High Risk: 51% - 100%\n\n"
        "Client Response Rule:\n"
        f"TBD is counted under Yes Risk. Current Yes Risk = {summary['yes_risk_count']} "
        f"({summary['yes_risk_percent']}%)."
    )

    for p in body.text_frame.paragraphs:
        p.font.size = Pt(16)


def add_detail_slides(prs, results, summary):
    rows_per_slide = 4

    for i in range(0, len(results), rows_per_slide):
        chunk = results[i:i + rows_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.6), Inches(0.4))
        header.text_frame.text = "Detailed Gold Standard Comparison"
        header.text_frame.paragraphs[0].font.size = Pt(20)
        header.text_frame.paragraphs[0].font.bold = True

        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12.6), Inches(0.3))
        sub.text_frame.text = (
            f"Overall Risk: {summary['overall_risk_level']} "
            f"({summary['comparison_risk_percent']}%) | "
            f"Low: {summary['low']} | Medium: {summary['medium']} | High: {summary['high']}"
        )
        sub.text_frame.paragraphs[0].font.size = Pt(11)

        table = slide.shapes.add_table(
            rows=len(chunk) + 1,
            cols=6,
            left=Inches(0.25),
            top=Inches(1.05),
            width=Inches(12.85),
            height=Inches(6.1)
        ).table

        widths = [1.45, 2.25, 3.1, 3.1, 1.0, 0.95]
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)

        headers = [
            "Domain",
            "Client Answer",
            "Gold Standard Answer",
            "Question",
            "Score",
            "Risk"
        ]

        for col, header_text in enumerate(headers):
            cell = table.cell(0, col)
            set_cell_text(cell, header_text, font_size=10, bold=True, font_color=RGBColor(255, 255, 255))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(161, 0, 255)

        for row_idx, item in enumerate(chunk, start=1):
            row_values = [
                item["domain"],
                item["client_answer"],
                item["gold_answer"],
                item["question"],
                f"{item['match_score']}%",
                item["risk_level"]
            ]

            for col_idx, value in enumerate(row_values):
                set_cell_text(table.cell(row_idx, col_idx), value, font_size=7)


def create_powerpoint(results, summary):
    os.makedirs(REPORT_FOLDER, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_summary_slide(prs, summary)
    add_logic_slide(prs, summary)
    add_detail_slides(prs, results, summary)

    prs.save(PPT_FILE_PATH)


@app.route("/", methods=["GET", "POST"])
def index():
    results = []
    summary = None
    error = None
    ppt_file = None

    if request.method == "POST":
        uploaded_file = request.files.get("client_file")

        if not uploaded_file:
            error = "Please select a client Excel file."
            return render_template(
                "compare.html",
                results=results,
                summary=summary,
                error=error,
                ppt_file=ppt_file
            )

        if not os.path.exists(GOLD_FILE):
            error = f"Gold standard file not found: {GOLD_FILE}"
            return render_template(
                "compare.html",
                results=results,
                summary=summary,
                error=error,
                ppt_file=ppt_file
            )

        try:
            client_df = load_client_file(uploaded_file)
            gold_map = load_gold_file()

            for _, row in client_df.iterrows():
                domain = row["Domain"]
                question = row["Question"]
                client_answer = row["Answer"]
                selected_risk = row["SelectedRisk"]

                gold_answer = gold_map.get(normalize_domain(domain), "")

                risk_level, match_score = compare_answers(client_answer, gold_answer)

                results.append({
                    "domain": domain,
                    "question": question,
                    "client_answer": client_answer,
                    "gold_answer": gold_answer,
                    "selected_risk": selected_risk,
                    "match_score": match_score,
                    "risk_level": risk_level
                })

            summary = calculate_summary(results)
            create_powerpoint(results, summary)
            ppt_file = PPT_FILE_NAME

        except Exception as e:
            error = str(e)

    return render_template(
        "compare.html",
        results=results,
        summary=summary,
        error=error,
        ppt_file=ppt_file
    )


if __name__ == "__main__":
    app.run(debug=True, port=5003)
