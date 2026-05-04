from flask import Flask, render_template, request, send_file
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os

app = Flask(__name__)

EXCEL_FILE = "mythos.xlsx"
PPT_FILE = "mythos_assessment_report.pptx"

QUESTIONS = [
    {"domain": "Identity inventory", "question": "Do you have a complete inventory of all human, service, machine, workload, and AI-agent identities?"},
    {"domain": "AI agent identity", "question": "Can you distinguish between a human user, automation script, service account, and AI agent in logs?"},
    {"domain": "Authentication", "question": "Which IdP is used: Entra ID, Okta, Ping, Oracle IAM, or others?"},
    {"domain": "MFA", "question": "Is MFA enforced for all privileged users, admins, contractors, and service owners?"},
    {"domain": "Conditional Access", "question": "Do you enforce device, location, risk, session, and workload-based access policies?"},
    {"domain": "Privileged access", "question": "Do admins use permanent access or just-in-time access through PIM/PAM?"},
    {"domain": "Service accounts", "question": "How many service accounts exist, and are any using static passwords or long-lived secrets?"},
    {"domain": "API access", "question": "Are APIs protected by OAuth, mTLS, API keys, APIM, gateway policies, or network restrictions?"},
    {"domain": "Token controls", "question": "Are JWTs validated for issuer, audience, expiry, scopes, roles, and signature?"},
    {"domain": "Secrets management", "question": "Are API keys, passwords, certificates, and model credentials stored in Key Vault or a secrets manager?"},
    {"domain": "Agent permissions", "question": "What systems can AI agents access: email, HR, CRM, code repos, cloud, tickets, databases, or identity systems?"},
    {"domain": "Least privilege", "question": "Are AI agents restricted to minimum required permissions?"},
    {"domain": "Human approval", "question": "Which AI-agent actions require human approval before execution?"},
    {"domain": "Tool access", "question": "Are AI agents allowed to call external tools, APIs, browsers, shells, or code execution environments?"},
    {"domain": "Policy enforcement", "question": "Do you use OPA, PBAC, ABAC, IAM policies, or workflow approval gates before agent actions?"},
    {"domain": "SaaS identity", "question": "Which SaaS platforms are connected to SSO, and which still use local accounts?"},
    {"domain": "Joiner-mover-leaver", "question": "How quickly are user and contractor accounts disabled after role change or exit?"},
    {"domain": "Identity governance", "question": "Do you run access reviews, SoD checks, entitlement reviews, and orphan-account cleanup?"},
    {"domain": "Logging", "question": "Do you log identity events from IdP, PAM, cloud, SaaS, APIM, and AI-agent platforms?"},
    {"domain": "Detection", "question": "Can you detect abnormal login, impossible travel, token replay, API abuse, or high-speed scanning?"},
    {"domain": "AI-agent behavior", "question": "Can you detect non-human behavior such as rapid enumeration, repeated API probing, or automated privilege testing?"},
    {"domain": "Model inventory", "question": "Do you know which AI models are deployed, where they run, and who can access them?"},
    {"domain": "Shadow AI", "question": "Can users deploy or access unauthorized AI models or agents?"},
    {"domain": "Network segmentation", "question": "Are identity systems, APIs, databases, and AI platforms microsegmented?"},
    {"domain": "Data access", "question": "Are retrieval systems protected with tenant, role, sensitivity, and data-label filters?"},
    {"domain": "Incident response", "question": "Do you have a playbook for compromised service account, AI agent misuse, or rogue model access?"},
    {"domain": "Mythos readiness", "question": "If an AI model finds identity vulnerabilities rapidly, can you prioritize, contain, and remediate before exploitation?"},
    {"domain": "Compensating controls", "question": "Where patching is delayed, do you have containment controls such as MFA, CA, PAM, segmentation, APIM throttling, and monitoring?"},
    {"domain": "Third-party access", "question": "Do vendors, contractors, and external AI tools have restricted and monitored access?"},
    {"domain": "Evidence", "question": "Can you produce audit evidence showing who accessed what, through which identity, and whether it was human or AI-driven?"}
]


def create_excel():
    """Create Excel file with the latest expected columns if it does not already exist."""
    if not os.path.exists(EXCEL_FILE):
        wb = Workbook()
        ws = wb.active
        ws.title = "Assessment"
        ws.append([
            "Timestamp", "Domain", "Question", "Answer", "Risk",
            "Yes %", "No %", "TBD %", "Overall Risk %", "Risk Level"
        ])
        wb.save(EXCEL_FILE)


def calculate_percentages(yes_count, no_count, tbd_count):
    """
    Calculate separate Yes/No/TBD percentages.

    Important business rule requested:
    TBD is counted under Yes Risk for overall risk scoring.

    Overall Risk % = (Yes + TBD) / Total * 100
    """
    total = yes_count + no_count + tbd_count

    if total == 0:
        return {
            "total": 0,
            "yes_pct": 0,
            "no_pct": 0,
            "tbd_pct": 0,
            "risk_percentage": 0,
            "risk_level": "Not Assessed"
        }

    yes_pct = round((yes_count / total) * 100, 2)
    no_pct = round((no_count / total) * 100, 2)
    tbd_pct = round((tbd_count / total) * 100, 2)

    # TBD is treated as Yes Risk based on your requirement.
    risk_yes_count = yes_count + tbd_count
    risk_percentage = round((risk_yes_count / total) * 100, 2)

    if risk_percentage <= 20:
        risk_level = "Low Risk"
    elif risk_percentage <= 50:
        risk_level = "Medium Risk"
    else:
        risk_level = "High Risk"

    return {
        "total": total,
        "yes_pct": yes_pct,
        "no_pct": no_pct,
        "tbd_pct": tbd_pct,
        "risk_percentage": risk_percentage,
        "risk_level": risk_level
    }


def risk_color(risk_level):
    if risk_level == "Low Risk":
        return RGBColor(34, 197, 94)      # green
    if risk_level == "Medium Risk":
        return RGBColor(245, 158, 11)     # amber
    if risk_level == "High Risk":
        return RGBColor(239, 68, 68)      # red
    return RGBColor(107, 114, 128)        # gray


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def create_powerpoint(rows, summary):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    yes_count = summary["yes_count"]
    no_count = summary["no_count"]
    tbd_count = summary["tbd_count"]
    total = summary["total"]
    yes_pct = summary["yes_pct"]
    no_pct = summary["no_pct"]
    tbd_pct = summary["tbd_pct"]
    risk_percentage = summary["risk_percentage"]
    risk_level = summary["risk_level"]

    # Slide 1: Executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.2), Inches(0.5),
                "Mythos Identity Assessment Report", 28, True, RGBColor(255, 255, 255))

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 10)

    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12.2), Inches(0.35),
                "Risk scoring rule: TBD is counted as Yes Risk for overall Mythos risk.",
                14, False, RGBColor(229, 231, 235))

    # Risk box
    risk_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.65), Inches(12.2), Inches(1.0))
    risk_box.fill.solid()
    risk_box.fill.fore_color.rgb = risk_color(risk_level)
    risk_box.line.color.rgb = risk_color(risk_level)
    risk_text = risk_box.text_frame
    risk_text.clear()
    p = risk_text.paragraphs[0]
    p.text = f"Overall Risk: {risk_level} | Risk Score: {risk_percentage}%"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Summary table
    table = slide.shapes.add_table(
        rows=5,
        cols=3,
        left=Inches(1.2),
        top=Inches(3.0),
        width=Inches(10.8),
        height=Inches(2.6)
    ).table

    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = Inches(3.6)

    data = [
        ["Metric", "Count", "Percentage"],
        ["Yes Risk", str(yes_count), f"{yes_pct}%"],
        ["No Risk", str(no_count), f"{no_pct}%"],
        ["TBD", str(tbd_count), f"{tbd_pct}%"],
        ["Total Questions", str(total), "100%"]
    ]

    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(161, 0, 255)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(0.8), Inches(6.25), Inches(12), Inches(0.4),
                "Risk thresholds: Low = 0-20%, Medium = 21-50%, High = 51-100%.",
                12, False, RGBColor(229, 231, 235), PP_ALIGN.CENTER)

    # Slide 2: Risk logic
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.5),
                "Risk Calculation Logic", 24, True, RGBColor(0, 0, 0))

    logic = (
        "1. Count each selected Risk value: Yes, No, and TBD.\n"
        "2. Calculate separate percentages for Yes / No / TBD.\n"
        "3. Treat TBD as Yes Risk for overall scoring.\n"
        "4. Overall Risk % = ((Yes Count + TBD Count) / Total Questions) × 100.\n"
        "5. Low Risk = 0-20%, Medium Risk = 21-50%, High Risk = 51-100%."
    )
    add_textbox(slide, Inches(0.7), Inches(1.2), Inches(12), Inches(2.0), logic, 18, False, RGBColor(0, 0, 0))

    # Slide 3+: Detailed results
    rows_per_slide = 4

    for i in range(0, len(rows), rows_per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.4),
                    "Detailed Assessment Results", 20, True, RGBColor(0, 0, 0))
        add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.3),
                    f"Overall Risk: {risk_level} | Risk Score: {risk_percentage}% | Yes: {yes_count}, No: {no_count}, TBD: {tbd_count}",
                    12, False, RGBColor(55, 65, 81))

        chunk = rows[i:i + rows_per_slide]

        table = slide.shapes.add_table(
            rows=len(chunk) + 1,
            cols=4,
            left=Inches(0.35),
            top=Inches(1.1),
            width=Inches(12.6),
            height=Inches(5.8)
        ).table

        table.columns[0].width = Inches(1.8)
        table.columns[1].width = Inches(4.0)
        table.columns[2].width = Inches(5.2)
        table.columns[3].width = Inches(1.6)

        headers = ["Domain", "Question", "Client Answer", "Risk"]
        for col, h in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(161, 0, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        for r_idx, item in enumerate(chunk, start=1):
            values = [item["domain"], item["question"], item["answer"], item["risk"]]
            for c_idx, val in enumerate(values):
                cell = table.cell(r_idx, c_idx)
                cell.text = val or "-"
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(0, 0, 0)

                if c_idx == 3:
                    cell.fill.solid()
                    if val == "Yes":
                        cell.fill.fore_color.rgb = RGBColor(254, 226, 226)
                    elif val == "No":
                        cell.fill.fore_color.rgb = RGBColor(220, 252, 231)
                    elif val == "TBD":
                        cell.fill.fore_color.rgb = RGBColor(254, 243, 199)

    prs.save(PPT_FILE)


@app.route("/", methods=["GET", "POST"])
def index():
    message = ""
    summary = None
    ppt_ready = False

    if request.method == "POST":
        create_excel()

        wb = load_workbook(EXCEL_FILE)
        ws = wb["Assessment"]

        rows = []
        yes_count = 0
        no_count = 0
        tbd_count = 0

        for item in QUESTIONS:
            domain = item["domain"]
            question = item["question"]
            answer = request.form.get(domain, "").strip()
            risk = request.form.get(domain + "_risk", "")

            if risk == "Yes":
                yes_count += 1
            elif risk == "No":
                no_count += 1
            elif risk == "TBD":
                tbd_count += 1

            rows.append({
                "domain": domain,
                "question": question,
                "answer": answer,
                "risk": risk
            })

        calc = calculate_percentages(yes_count, no_count, tbd_count)

        summary = {
            "yes_count": yes_count,
            "no_count": no_count,
            "tbd_count": tbd_count,
            **calc
        }

        for item in rows:
            ws.append([
                datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                item["domain"],
                item["question"],
                item["answer"],
                item["risk"],
                summary["yes_pct"],
                summary["no_pct"],
                summary["tbd_pct"],
                summary["risk_percentage"],
                summary["risk_level"]
            ])

        wb.save(EXCEL_FILE)
        create_powerpoint(rows, summary)

        ppt_ready = True
        message = (
            f"Saved. PPT generated. Overall Risk: {summary['risk_level']} "
            f"({summary['risk_percentage']}%). "
            f"Yes: {summary['yes_count']} ({summary['yes_pct']}%), "
            f"No: {summary['no_count']} ({summary['no_pct']}%), "
            f"TBD: {summary['tbd_count']} ({summary['tbd_pct']}%). "
            "TBD is counted under Yes Risk."
        )

    return render_template(
        "index.html",
        questions=QUESTIONS,
        message=message,
        summary=summary,
        ppt_ready=ppt_ready
    )


@app.route("/download-ppt")
def download_ppt():
    if not os.path.exists(PPT_FILE):
        return "PPT report has not been generated yet. Submit the assessment first.", 404

    return send_file(PPT_FILE, as_attachment=True)


if __name__ == "__main__":
    create_excel()
    app.run(debug=True, port=5001)
